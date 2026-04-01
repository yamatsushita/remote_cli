"""
Generate a PowerPoint presentation explaining Structure-from-Motion (SfM)
and why scale ambiguity remains in monocular SfM solutions.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ---------------------------------------------------------------------------
# Colour palette
# ---------------------------------------------------------------------------
BG_DARK = RGBColor(0x1B, 0x1B, 0x2F)       # dark navy background
ACCENT_BLUE = RGBColor(0x00, 0x9B, 0xE6)    # bright blue accent
ACCENT_ORANGE = RGBColor(0xF4, 0x84, 0x2D)  # orange accent
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
DARK_TEXT = RGBColor(0x22, 0x22, 0x22)
SOFT_BG = RGBColor(0xF0, 0xF4, 0xF8)        # light background for content slides

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


def set_slide_bg(slide, color):
    """Set a solid background colour for a slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, *,
                font_size=18, bold=False, color=DARK_TEXT,
                alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    """Helper to add a simple text box."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_slide(slide, items, left, top, width, height, *,
                     font_size=18, color=DARK_TEXT, spacing=Pt(8)):
    """Add a bulleted text frame to the slide."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = spacing
        p.level = 0
    return txBox


def draw_accent_bar(slide, left, top, width, height, color=ACCENT_BLUE):
    """Draw a small coloured rectangle as an accent bar."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


# ===================================================================
# SLIDE BUILDERS
# ===================================================================

def slide_title(prs):
    """Slide 1 – Title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide, BG_DARK)

    draw_accent_bar(slide, Inches(0), Inches(3.1), SLIDE_WIDTH, Inches(0.06), ACCENT_BLUE)

    add_textbox(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.5),
                "Structure-from-Motion",
                font_size=48, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

    add_textbox(slide, Inches(1), Inches(3.4), Inches(11), Inches(1.0),
                "Understanding Scale Ambiguity in Monocular 3-D Reconstruction",
                font_size=24, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

    add_textbox(slide, Inches(1), Inches(5.5), Inches(11), Inches(0.6),
                "Computer Vision  •  Multi-View Geometry  •  3-D Reconstruction",
                font_size=16, color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER)


def slide_overview(prs):
    """Slide 2 – What is SfM?"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, SOFT_BG)
    draw_accent_bar(slide, Inches(0.8), Inches(0.7), Inches(0.08), Inches(0.55), ACCENT_BLUE)
    add_textbox(slide, Inches(1.1), Inches(0.6), Inches(10), Inches(0.7),
                "What is Structure-from-Motion?",
                font_size=32, bold=True, color=DARK_TEXT)

    items = [
        "• Structure-from-Motion (SfM) recovers 3-D structure and camera motion "
        "from a set of 2-D images taken from different viewpoints.",
        "• It is a foundational technique in computer vision, used in autonomous "
        "driving, augmented reality, and cultural-heritage digitisation.",
        "• The inputs are 2-D feature correspondences across images; the outputs "
        "are a 3-D point cloud and the camera poses (rotation + translation).",
        "• SfM solves an inverse problem: projecting a 3-D scene onto 2-D images "
        "loses one dimension of information — depth.",
    ]
    add_bullet_slide(slide, items,
                     Inches(1), Inches(1.8), Inches(11), Inches(5),
                     font_size=20, color=DARK_TEXT)


def slide_pipeline(prs):
    """Slide 3 – The SfM Pipeline."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, SOFT_BG)
    draw_accent_bar(slide, Inches(0.8), Inches(0.7), Inches(0.08), Inches(0.55), ACCENT_BLUE)
    add_textbox(slide, Inches(1.1), Inches(0.6), Inches(10), Inches(0.7),
                "The SfM Pipeline",
                font_size=32, bold=True, color=DARK_TEXT)

    steps = [
        ("1", "Feature Detection & Matching",
         "Detect key-points (e.g. SIFT, ORB) in each image and match them across views."),
        ("2", "Essential / Fundamental Matrix Estimation",
         "Compute the geometric relationship between two views using matched points."),
        ("3", "Camera Pose Recovery",
         "Decompose the Essential Matrix into rotation R and translation t (up to scale)."),
        ("4", "Triangulation",
         "Back-project matched 2-D points into 3-D using the recovered poses."),
        ("5", "Bundle Adjustment",
         "Jointly refine all 3-D points and camera parameters to minimise re-projection error."),
    ]

    y = Inches(1.8)
    for num, title, desc in steps:
        # Step number circle
        shape = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                       Inches(1.0), y, Inches(0.55), Inches(0.55))
        shape.fill.solid()
        shape.fill.fore_color.rgb = ACCENT_BLUE
        shape.line.fill.background()
        tf = shape.text_frame
        tf.paragraphs[0].text = num
        tf.paragraphs[0].font.size = Pt(18)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = WHITE
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE

        add_textbox(slide, Inches(1.8), y - Inches(0.05), Inches(10), Inches(0.4),
                    title, font_size=20, bold=True, color=DARK_TEXT)
        add_textbox(slide, Inches(1.8), y + Inches(0.35), Inches(10), Inches(0.55),
                    desc, font_size=16, color=RGBColor(0x55, 0x55, 0x55))
        y += Inches(1.05)


def slide_epipolar(prs):
    """Slide 4 – Epipolar Geometry & the Essential Matrix."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, SOFT_BG)
    draw_accent_bar(slide, Inches(0.8), Inches(0.7), Inches(0.08), Inches(0.55), ACCENT_BLUE)
    add_textbox(slide, Inches(1.1), Inches(0.6), Inches(10), Inches(0.7),
                "Epipolar Geometry & the Essential Matrix",
                font_size=32, bold=True, color=DARK_TEXT)

    items = [
        "• For calibrated cameras the epipolar constraint is:",
        "          x'ᵀ  E  x  =  0",
        "  where x, x' are normalised image coordinates and E is the Essential Matrix.",
        "",
        "• The Essential Matrix encodes the relative pose:",
        "          E  =  [t]×  R",
        "  where R is the rotation and  [t]× is the skew-symmetric matrix of the "
        "translation vector t.",
        "",
        "• E has five degrees of freedom (3 for R, 2 for the direction of t).",
        "  Note: the magnitude (norm) of t is NOT recoverable from E.",
    ]
    add_bullet_slide(slide, items,
                     Inches(1), Inches(1.7), Inches(11), Inches(5),
                     font_size=19, color=DARK_TEXT, spacing=Pt(4))


def slide_scale_ambiguity_intro(prs):
    """Slide 5 – The Scale Ambiguity Problem."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    draw_accent_bar(slide, Inches(0.8), Inches(0.7), Inches(0.08), Inches(0.55), ACCENT_ORANGE)
    add_textbox(slide, Inches(1.1), Inches(0.6), Inches(10), Inches(0.7),
                "The Scale Ambiguity Problem",
                font_size=32, bold=True, color=WHITE)

    items = [
        "• Monocular SfM can recover the shape of the 3-D scene and the direction "
        "of camera motion, but NOT the absolute size (scale) of the reconstruction.",
        "",
        "• A toy house photographed close-up and a real house photographed from "
        "far away can produce identical images.",
        "",
        "• Mathematically, if  (R, t, {Xᵢ})  is a valid solution, then  "
        "(R, αt, {αXᵢ})  for any α > 0  is equally valid.",
        "",
        "• This means the reconstruction is determined only up to a global scale "
        "factor — this is known as scale ambiguity.",
    ]
    add_bullet_slide(slide, items,
                     Inches(1), Inches(1.7), Inches(11), Inches(5),
                     font_size=20, color=LIGHT_GRAY, spacing=Pt(4))


def slide_why_scale_ambiguity(prs):
    """Slide 6 – Why Does Scale Ambiguity Exist? (Mathematical Reason)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, SOFT_BG)
    draw_accent_bar(slide, Inches(0.8), Inches(0.7), Inches(0.08), Inches(0.55), ACCENT_ORANGE)
    add_textbox(slide, Inches(1.1), Inches(0.6), Inches(10), Inches(0.7),
                "Why Does Scale Ambiguity Exist?",
                font_size=32, bold=True, color=DARK_TEXT)

    items = [
        "1.  Projection eliminates depth — The pinhole camera model projects a "
        "3-D point  X  to  x = K [R | t] X.  Multiplying  X  and  t  by "
        "the same factor  α  leaves  x  unchanged because the perspective "
        "division cancels α.",
        "",
        "2.  The Essential Matrix is scale-blind — E = [t]× R  is unchanged "
        "when  t  is replaced by  αt  (the cross-product matrix [αt]× = α[t]×, "
        "and E is defined only up to scale).  The 5-point / 8-point algorithms "
        "therefore recover  t  only as a unit vector.",
        "",
        "3.  Triangulation inherits the ambiguity — Triangulated 3-D points are "
        "computed from the recovered  (R, t).  Since  ‖t‖  is unknown, the "
        "triangulated depth is only determined up to the same unknown factor.",
        "",
        "4.  Fundamental reason — A single camera provides bearing (direction) "
        "measurements, not range (distance) measurements.  Without at least "
        "one known real-world distance, no number of images from a monocular "
        "camera can determine the absolute scale.",
    ]
    add_bullet_slide(slide, items,
                     Inches(1), Inches(1.7), Inches(11.5), Inches(5.2),
                     font_size=17, color=DARK_TEXT, spacing=Pt(6))


def slide_geometric_intuition(prs):
    """Slide 7 – Geometric Intuition."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, SOFT_BG)
    draw_accent_bar(slide, Inches(0.8), Inches(0.7), Inches(0.08), Inches(0.55), ACCENT_ORANGE)
    add_textbox(slide, Inches(1.1), Inches(0.6), Inches(10), Inches(0.7),
                "Geometric Intuition for Scale Ambiguity",
                font_size=32, bold=True, color=DARK_TEXT)

    # Left column – small scene
    add_textbox(slide, Inches(1.5), Inches(1.7), Inches(5), Inches(0.5),
                "Scenario A – Small scene, small baseline",
                font_size=20, bold=True, color=ACCENT_BLUE)

    add_textbox(slide, Inches(1.5), Inches(2.3), Inches(5), Inches(1.6),
                "A tabletop object at distance d, with the camera "
                "moving by baseline b, produces a set of 2-D images.",
                font_size=17, color=DARK_TEXT)

    # Diagram: two small cameras viewing a small cube
    _draw_camera_icon(slide, Inches(2.0), Inches(4.0), ACCENT_BLUE, "C₁")
    _draw_camera_icon(slide, Inches(4.5), Inches(4.0), ACCENT_BLUE, "C₂")
    _draw_cube(slide, Inches(3.1), Inches(5.2), Inches(0.6), ACCENT_BLUE)

    add_textbox(slide, Inches(2.5), Inches(4.6), Inches(2.5), Inches(0.4),
                "baseline b", font_size=14, color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER)

    # Right column – large scene
    add_textbox(slide, Inches(7.5), Inches(1.7), Inches(5), Inches(0.5),
                "Scenario B – Large scene, large baseline",
                font_size=20, bold=True, color=ACCENT_ORANGE)

    add_textbox(slide, Inches(7.5), Inches(2.3), Inches(5), Inches(1.6),
                "A building at distance αd, with the camera moving "
                "by baseline αb, produces the SAME 2-D images.",
                font_size=17, color=DARK_TEXT)

    _draw_camera_icon(slide, Inches(7.8), Inches(4.0), ACCENT_ORANGE, "C₁")
    _draw_camera_icon(slide, Inches(10.8), Inches(4.0), ACCENT_ORANGE, "C₂")
    _draw_cube(slide, Inches(9.0), Inches(5.0), Inches(1.0), ACCENT_ORANGE)

    add_textbox(slide, Inches(8.5), Inches(4.6), Inches(3.0), Inches(0.4),
                "baseline αb", font_size=14, color=ACCENT_ORANGE, alignment=PP_ALIGN.CENTER)

    # Bottom note
    add_textbox(slide, Inches(1.5), Inches(6.4), Inches(10.5), Inches(0.7),
                "Both scenarios yield identical 2-D projections → SfM cannot "
                "distinguish between them without external scale information.",
                font_size=18, bold=True, color=DARK_TEXT, alignment=PP_ALIGN.CENTER)


def _draw_camera_icon(slide, left, top, color, label):
    """Draw a small trapezoid-like camera icon using a rectangle + triangle."""
    # Camera body
    body = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  left, top, Inches(0.5), Inches(0.35))
    body.fill.solid()
    body.fill.fore_color.rgb = color
    body.line.fill.background()
    # Label
    add_textbox(slide, left, top - Inches(0.35), Inches(0.5), Inches(0.35),
                label, font_size=13, bold=True, color=color, alignment=PP_ALIGN.CENTER)


def _draw_cube(slide, left, top, size, color):
    """Draw a simple square to represent the 3-D object."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.color.rgb = color
    shape.fill.fore_color.brightness = 0.6


def slide_math_proof(prs):
    """Slide 8 – Formal Proof of Scale Ambiguity."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, SOFT_BG)
    draw_accent_bar(slide, Inches(0.8), Inches(0.7), Inches(0.08), Inches(0.55), ACCENT_ORANGE)
    add_textbox(slide, Inches(1.1), Inches(0.6), Inches(10), Inches(0.7),
                "Formal Argument: Scale is Unobservable",
                font_size=32, bold=True, color=DARK_TEXT)

    lines = [
        "Consider a 3-D point  X  viewed by two cameras with poses  (I, 0)  and  (R, t).",
        "",
        "The projection in the second camera is:",
        "       x'  ∝  K ( R X + t )",
        "",
        "Now replace  t → αt  and  X → αX  (scale the world and the baseline equally):",
        "       x'  ∝  K ( R (αX) + αt )  =  K · α ( R X + t )  ∝  K ( R X + t )",
        "",
        "The factor  α  cancels in the projective (homogeneous) division.",
        "",
        "Therefore the 2-D observations are invariant to a global scaling of the scene "
        "and translation — the scale factor  α  is unobservable from images alone.",
        "",
        "Key insight:  Perspective projection is a  ℝ³ → ℝ²  mapping that preserves "
        "angles (bearings) but not distances.  Scale lives in the  null space  of the "
        "measurement model.",
    ]
    add_bullet_slide(slide, lines,
                     Inches(1), Inches(1.6), Inches(11.5), Inches(5.5),
                     font_size=17, color=DARK_TEXT, spacing=Pt(3))


def slide_resolving(prs):
    """Slide 9 – How to Resolve Scale Ambiguity."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, SOFT_BG)
    draw_accent_bar(slide, Inches(0.8), Inches(0.7), Inches(0.08), Inches(0.55), ACCENT_BLUE)
    add_textbox(slide, Inches(1.1), Inches(0.6), Inches(10), Inches(0.7),
                "Resolving Scale Ambiguity in Practice",
                font_size=32, bold=True, color=DARK_TEXT)

    methods = [
        ("Known Object Size",
         "Place an object of known dimension in the scene (e.g. a calibration "
         "target).  One real-world measurement pins down the global scale."),
        ("Stereo / Multi-Camera Rig",
         "A calibrated stereo pair with a known baseline provides absolute scale "
         "because the inter-camera distance is fixed and known."),
        ("IMU / Inertial Sensors",
         "An accelerometer measures metric acceleration; fusing IMU data with "
         "visual SfM (Visual-Inertial Odometry) recovers metric scale."),
        ("GPS / GNSS",
         "Geo-tagged images supply metric positions that anchor the scale."),
        ("Deep Monocular Depth",
         "Learned depth-prediction networks can provide a metric depth prior, "
         "though accuracy depends on training data."),
    ]

    y = Inches(1.7)
    for title, desc in methods:
        add_textbox(slide, Inches(1.2), y, Inches(3.5), Inches(0.4),
                    "▸  " + title, font_size=19, bold=True, color=ACCENT_BLUE)
        add_textbox(slide, Inches(4.8), y, Inches(8), Inches(0.75),
                    desc, font_size=16, color=DARK_TEXT)
        y += Inches(0.95)


def slide_summary(prs):
    """Slide 10 – Summary / Take-aways."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    draw_accent_bar(slide, Inches(0.8), Inches(0.7), Inches(0.08), Inches(0.55), ACCENT_BLUE)
    add_textbox(slide, Inches(1.1), Inches(0.6), Inches(10), Inches(0.7),
                "Key Take-aways",
                font_size=32, bold=True, color=WHITE)

    items = [
        "1.  SfM reconstructs 3-D geometry and camera motion from 2-D images.",
        "",
        "2.  The Essential Matrix encodes relative pose, but only the direction "
        "of translation — not its magnitude.",
        "",
        "3.  Because perspective projection divides out depth, uniformly scaling "
        "the scene and the baseline produces identical images "
        "→  scale ambiguity.",
        "",
        "4.  Scale ambiguity is intrinsic to monocular vision; it lives in the "
        "null space of the projection.",
        "",
        "5.  In practice, scale is resolved by injecting at least one metric "
        "measurement: a known distance, a stereo baseline, or an IMU.",
    ]
    add_bullet_slide(slide, items,
                     Inches(1), Inches(1.6), Inches(11), Inches(5.5),
                     font_size=19, color=LIGHT_GRAY, spacing=Pt(4))


def slide_references(prs):
    """Slide 11 – References."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, SOFT_BG)
    draw_accent_bar(slide, Inches(0.8), Inches(0.7), Inches(0.08), Inches(0.55), ACCENT_BLUE)
    add_textbox(slide, Inches(1.1), Inches(0.6), Inches(10), Inches(0.7),
                "References",
                font_size=32, bold=True, color=DARK_TEXT)

    refs = [
        "• Hartley, R. & Zisserman, A. (2004). Multiple View Geometry in "
        "Computer Vision, 2nd ed. Cambridge University Press.",
        "",
        "• Szeliski, R. (2022). Computer Vision: Algorithms and Applications, "
        "2nd ed. Springer.",
        "",
        "• Nistér, D. (2004). An efficient solution to the five-point relative "
        "pose problem. IEEE TPAMI, 26(6), 756–770.",
        "",
        "• Scaramuzza, D. & Fraundorfer, F. (2011). Visual Odometry — Part I & II. "
        "IEEE Robotics & Automation Magazine.",
    ]
    add_bullet_slide(slide, refs,
                     Inches(1), Inches(1.7), Inches(11), Inches(5),
                     font_size=17, color=DARK_TEXT, spacing=Pt(3))


# ===================================================================
# MAIN
# ===================================================================

def main():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    slide_title(prs)
    slide_overview(prs)
    slide_pipeline(prs)
    slide_epipolar(prs)
    slide_scale_ambiguity_intro(prs)
    slide_why_scale_ambiguity(prs)
    slide_geometric_intuition(prs)
    slide_math_proof(prs)
    slide_resolving(prs)
    slide_summary(prs)
    slide_references(prs)

    out = "sfm_scale_ambiguity.pptx"
    prs.save(out)
    print(f"Presentation saved to {out}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
