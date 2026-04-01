#!/usr/bin/env python3
"""Generate a PowerPoint presentation explaining Structure-from-Motion
and the scale ambiguity problem."""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# -- Color palette --
BG_DARK = RGBColor(0x1B, 0x1B, 0x2F)
BG_MID = RGBColor(0x24, 0x24, 0x3E)
ACCENT = RGBColor(0x00, 0xD2, 0xFF)
ACCENT2 = RGBColor(0xFF, 0x6B, 0x6B)
ACCENT3 = RGBColor(0x4E, 0xCB, 0x71)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xBB, 0xBB, 0xCC)
YELLOW = RGBColor(0xFF, 0xD9, 0x3D)


def set_slide_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text,
                font_size=18, color=WHITE, bold=False,
                alignment=PP_ALIGN.LEFT, font_name="Segoe UI"):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tf


def add_bullet_list(slide, left, top, width, height, items,
                    font_size=18, color=WHITE, spacing=Pt(8)):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Segoe UI"
        p.space_after = spacing
        pPr = p._p.get_or_add_pPr()
        for child in list(pPr):
            if "buChar" in child.tag or "buNone" in child.tag:
                pPr.remove(child)
        pPr.append(pPr.makeelement(qn("a:buChar"), {"char": "\u25B8"}))
    return tf


def add_rounded_rect(slide, left, top, width, height, fill_color,
                     text="", font_size=14, font_color=WHITE):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].text = text
        tf.paragraphs[0].font.size = Pt(font_size)
        tf.paragraphs[0].font.color.rgb = font_color
        tf.paragraphs[0].font.name = "Segoe UI"
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    return shape


def add_separator(slide, x, y, w):
    s = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Pt(3))
    s.fill.solid()
    s.fill.fore_color.rgb = ACCENT
    s.line.fill.background()


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ── SLIDE 1: Title ──────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_textbox(slide, 1, 1.5, 11.3, 1.2, "Structure-from-Motion",
                48, ACCENT, True, PP_ALIGN.CENTER)
    add_textbox(slide, 1, 2.8, 11.3, 0.8,
                "3D Reconstruction from 2D Images",
                28, WHITE, False, PP_ALIGN.CENTER)
    add_textbox(slide, 1, 4.0, 11.3, 0.6,
                "\u2014 and the Scale Ambiguity Problem \u2014",
                22, YELLOW, False, PP_ALIGN.CENTER)
    s = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(4), Inches(3.7), Inches(5.3), Pt(3))
    s.fill.solid(); s.fill.fore_color.rgb = ACCENT; s.line.fill.background()
    add_textbox(slide, 1, 5.5, 11.3, 0.5, "Computer Vision Fundamentals",
                16, LIGHT_GRAY, False, PP_ALIGN.CENTER)

    # ── SLIDE 2: What is SfM? ───────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_textbox(slide, 0.6, 0.3, 12, 0.8,
                "What is Structure-from-Motion?", 36, ACCENT, True)
    add_separator(slide, 0.6, 1.05, 5.0)
    add_textbox(slide, 0.6, 1.3, 6, 0.6,
                "Structure-from-Motion (SfM) is the process of recovering:",
                20, WHITE)
    add_bullet_list(slide, 0.8, 1.9, 5.8, 3.5, [
        "3D structure (point cloud) of a scene",
        "Camera poses (position + orientation) for each image",
        "Intrinsic camera parameters (focal length, etc.)",
    ], 19, WHITE)
    add_textbox(slide, 0.6, 3.6, 6, 1.0,
                "\u2026 all from a collection of 2D images taken from "
                "different viewpoints, without any prior knowledge of "
                "the scene geometry.", 18, LIGHT_GRAY)
    add_rounded_rect(slide, 7.2, 1.2, 5.5, 1.0, BG_MID,
                     "Input: Unordered 2D Images", 18, ACCENT)
    add_rounded_rect(slide, 7.2, 2.6, 5.5, 1.0, BG_MID,
                     "Process: Feature Matching \u2192 Geometry Estimation",
                     16, WHITE)
    add_rounded_rect(slide, 7.2, 4.0, 5.5, 1.0, BG_MID,
                     "Output: 3D Points + Camera Poses", 18, ACCENT3)
    for y in [2.15, 3.55]:
        a = slide.shapes.add_shape(
            MSO_SHAPE.DOWN_ARROW, Inches(9.7), Inches(y),
            Inches(0.4), Inches(0.4))
        a.fill.solid(); a.fill.fore_color.rgb = ACCENT
        a.line.fill.background()
    add_textbox(slide, 0.6, 5.0, 12, 0.8,
                "Applications:  autonomous driving  \u2022  augmented reality  "
                "\u2022  cultural heritage digitization  \u2022  robotics  "
                "\u2022  urban mapping", 16, LIGHT_GRAY)

    # ── SLIDE 3: SfM Pipeline ───────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_textbox(slide, 0.6, 0.3, 12, 0.8,
                "The SfM Pipeline", 36, ACCENT, True)
    add_separator(slide, 0.6, 1.05, 3.5)
    steps = [
        ("1", "Feature Detection\n& Description",
         "SIFT, ORB, SuperPoint\ndetect keypoints in\neach image"),
        ("2", "Feature Matching",
         "Match keypoints across\nimage pairs using\ndescriptor similarity"),
        ("3", "Geometric\nVerification",
         "Estimate fundamental /\nessential matrix;\nreject outliers (RANSAC)"),
        ("4", "Incremental\nReconstruction",
         "Initialize with 2-view\nreconstruction, then add\nimages one by one"),
        ("5", "Triangulation",
         "Compute 3D points from\ncorresponding 2D\nobservations"),
        ("6", "Bundle\nAdjustment",
         "Jointly optimize all\n3D points & camera\nparameters"),
    ]
    x0, bw, gap = 0.5, 1.95, 0.15
    for i, (num, title, desc) in enumerate(steps):
        x = x0 + i * (bw + gap)
        circ = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(x + 0.7), Inches(1.4),
            Inches(0.55), Inches(0.55))
        circ.fill.solid(); circ.fill.fore_color.rgb = ACCENT
        circ.line.fill.background()
        tf = circ.text_frame
        tf.paragraphs[0].text = num
        tf.paragraphs[0].font.size = Pt(18)
        tf.paragraphs[0].font.color.rgb = BG_DARK
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf.paragraphs[0].font.name = "Segoe UI"
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        add_rounded_rect(slide, x, 2.2, bw, 1.1, BG_MID, title, 14, WHITE)
        add_textbox(slide, x, 3.5, bw, 1.5, desc,
                    12, LIGHT_GRAY, False, PP_ALIGN.CENTER)
    for i in range(5):
        ax = x0 + (i + 1) * (bw + gap) - gap + 0.02
        ar = slide.shapes.add_shape(
            MSO_SHAPE.RIGHT_ARROW, Inches(ax - 0.15), Inches(2.6),
            Inches(0.25), Inches(0.2))
        ar.fill.solid(); ar.fill.fore_color.rgb = ACCENT
        ar.line.fill.background()
    add_textbox(slide, 0.6, 5.3, 12, 0.8,
                "Key insight: The entire pipeline works only with "
                "relative geometry \u2014 there is no absolute reference "
                "frame or scale.", 18, YELLOW, True)

    # ── SLIDE 4: Epipolar Geometry ──────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_textbox(slide, 0.6, 0.3, 12, 0.8,
                "Mathematical Foundation: Epipolar Geometry",
                36, ACCENT, True)
    add_separator(slide, 0.6, 1.05, 7.0)
    add_textbox(slide, 0.6, 1.3, 6, 0.6,
                "The Fundamental Matrix Constraint", 24, WHITE, True)
    add_textbox(slide, 0.6, 1.9, 6, 0.6,
                "For corresponding points x and x\u2019 in two images:",
                18, LIGHT_GRAY)
    add_rounded_rect(slide, 1.2, 2.5, 4.5, 0.9, BG_MID,
                     "x\u2019\u1d40 F x  =  0", 32, YELLOW)
    add_textbox(slide, 0.6, 3.7, 6, 0.5,
                "The Essential Matrix (calibrated case):",
                24, WHITE, True)
    add_rounded_rect(slide, 1.2, 4.3, 4.5, 0.9, BG_MID,
                     "E  =  [t]\u2093 R      (E = K\u2019\u1d40 F K)",
                     24, YELLOW)
    add_textbox(slide, 0.6, 5.5, 6, 1.0,
                "where  R = rotation,  t = translation,  "
                "[t]\u2093 = skew-symmetric matrix of t",
                16, LIGHT_GRAY)
    add_textbox(slide, 7.2, 1.3, 5.5, 0.6,
                "Key Properties of E = [t]\u2093 R", 22, ACCENT3, True)
    add_bullet_list(slide, 7.4, 1.9, 5.3, 4.5, [
        "E encodes the relative pose (R, t) between two cameras",
        "E has rank 2  (det(E) = 0)",
        "Decomposition of E \u2192 (R, t) is possible via SVD",
        "CRITICAL:  E = [t]\u2093 R  =  [\u03b1t]\u2093 R  for any scalar \u03b1 \u2260 0",
        "The translation t can only be recovered up to scale",
        "Scaling t by any positive factor \u03b1 gives the same E",
    ], 17, WHITE)
    add_textbox(slide, 7.2, 5.3, 5.5, 1.0,
                "\u2192  This is the mathematical root of scale ambiguity!",
                20, ACCENT2, True)

    # ── SLIDE 5: Why Scale Ambiguity Exists ─────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_textbox(slide, 0.6, 0.3, 12, 0.8,
                "Why Does Scale Ambiguity Exist?", 36, ACCENT, True)
    add_separator(slide, 0.6, 1.05, 6.0)
    add_textbox(slide, 0.6, 1.3, 12, 0.6,
                "The Core Insight: Projection Destroys Absolute Scale",
                26, YELLOW, True)
    add_rounded_rect(slide, 0.6, 2.2, 5.5, 4.2, BG_MID)
    add_textbox(slide, 0.8, 2.3, 5.1, 0.5,
                "Scenario A: Small Object, Close Cameras",
                18, ACCENT3, True)
    add_textbox(slide, 0.8, 2.9, 5.1, 2.5,
                "\u2022 A toy car (10 cm) photographed by two cameras "
                "20 cm apart\n\n"
                "\u2022 The images captured are identical to Scenario B\n\n"
                "\u2022 2D pixel coordinates: exactly the same\n\n"
                "\u2022 All geometric relationships preserved",
                16, WHITE)
    add_rounded_rect(slide, 7.0, 2.2, 5.7, 4.2, BG_MID)
    add_textbox(slide, 7.2, 2.3, 5.3, 0.5,
                "Scenario B: Large Object, Far Cameras",
                18, ACCENT2, True)
    add_textbox(slide, 7.2, 2.9, 5.3, 2.5,
                "\u2022 A real car (2 m) photographed by two cameras "
                "4 m apart\n\n"
                "\u2022 Scale factor = 20\u00d7  for both object and baseline\n\n"
                "\u2022 2D pixel coordinates: exactly the same\n\n"
                "\u2022 Impossible to distinguish from Scenario A using "
                "images alone",
                16, WHITE)
    add_textbox(slide, 5.9, 3.5, 1.3, 1.0, "=",
                60, YELLOW, True, PP_ALIGN.CENTER)
    add_textbox(slide, 0.6, 6.6, 12, 0.6,
                "Perspective projection is invariant to uniform scaling "
                "of the scene and camera positions.",
                18, YELLOW, True, PP_ALIGN.CENTER)

    # ── SLIDE 6: Mathematical Proof ─────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_textbox(slide, 0.6, 0.3, 12, 0.8,
                "Mathematical Proof of Scale Ambiguity",
                36, ACCENT, True)
    add_separator(slide, 0.6, 1.05, 6.5)
    add_textbox(slide, 0.6, 1.2, 12, 0.5,
                "Camera Projection Model", 24, WHITE, True)
    add_rounded_rect(slide, 0.6, 1.8, 5.8, 1.1, BG_MID,
                     "A 3D point X projects to image point x:\n"
                     " x  =  K [ R | t ] X", 17, WHITE)
    add_textbox(slide, 0.6, 3.1, 12, 0.5,
                "Now scale the scene and translations by \u03b1 > 0:",
                20, YELLOW, True)
    add_rounded_rect(slide, 0.6, 3.7, 12, 1.4, BG_MID,
                     "  X\u2019 = \u03b1X ,   t\u2019 = \u03b1t   "
                     "(scale everything by \u03b1)\n\n"
                     "  x  =  K [ R | t ] X   =   K [ R | \u03b1t ] "
                     "(\u03b1X)   =   K ( \u03b1RX + \u03b1t )  "
                     "=  \u03b1 \u00b7 K(RX + t)\n\n"
                     "  In homogeneous coordinates:  x  \u223c  "
                     "K(RX + t)  \u223c  K(\u03b1RX + \u03b1t)  "
                     "=  K [ R | \u03b1t ](\u03b1X)", 14, WHITE)
    add_textbox(slide, 0.6, 5.3, 12, 0.5,
                "Conclusion:", 22, ACCENT3, True)
    add_bullet_list(slide, 0.8, 5.8, 12, 1.5, [
        "The projection x is unchanged \u2014 the \u03b1 cancels in "
        "homogeneous coordinates",
        "SfM can recover (R, t) only up to a global scale factor on t",
        "Equivalently: the scene {X} and baselines {t} are recoverable "
        "only up to a shared scale",
    ], 17, WHITE)

    # ── SLIDE 7: Degrees of Freedom ─────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_textbox(slide, 0.6, 0.3, 12, 0.8,
                "Degrees of Freedom & the Gauge Ambiguity",
                36, ACCENT, True)
    add_separator(slide, 0.6, 1.05, 7.5)
    add_textbox(slide, 0.6, 1.3, 6, 0.5,
                "SfM recovers geometry up to a similarity transform:",
                20, WHITE)
    add_rounded_rect(slide, 0.6, 2.0, 5.8, 3.5, BG_MID)
    add_textbox(slide, 0.8, 2.1, 5.4, 0.5,
                "7 Degrees of Freedom (Gauge Freedom)",
                20, ACCENT, True)
    dofs = [
        ("3 DoF", "Global Rotation",
         "No absolute orientation reference", ACCENT3),
        ("3 DoF", "Global Translation",
         "No absolute position reference", ACCENT3),
        ("1 DoF", "Global Scale",
         "No absolute scale reference", ACCENT2),
    ]
    for i, (dof, name, desc, col) in enumerate(dofs):
        y = 2.7 + i * 0.9
        add_textbox(slide, 1.0, y, 1.2, 0.4, dof, 18, col, True)
        add_textbox(slide, 2.2, y, 2.0, 0.4, name, 18, WHITE, True)
        add_textbox(slide, 2.2, y + 0.35, 3.2, 0.4, desc, 14, LIGHT_GRAY)
    add_textbox(slide, 7.0, 1.3, 5.5, 0.5,
                "Why can\u2019t more images fix scale?",
                22, ACCENT2, True)
    add_bullet_list(slide, 7.2, 1.9, 5.3, 4.5, [
        "Adding more images adds more constraints",
        "But it also adds more unknowns (new camera pose)",
        "Scale is a global ambiguity \u2014 it cannot be resolved "
        "by adding more monocular images",
        "The fundamental / essential matrix only encodes direction "
        "of translation, not magnitude",
        "Triangulation gives depth ratios, not absolute depth",
        "The reprojection error (used in bundle adjustment) is "
        "also scale-invariant",
    ], 16, WHITE)
    add_textbox(slide, 0.6, 5.8, 12, 0.8,
                "Standard convention: fix the baseline of the first "
                "camera pair to ||t\u2081\u2082|| = 1  (unit translation).",
                18, YELLOW, True, PP_ALIGN.CENTER)

    # ── SLIDE 8: Visual Illustration ────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_textbox(slide, 0.6, 0.3, 12, 0.8,
                "Scale Ambiguity: Same Images, Different Worlds",
                36, ACCENT, True)
    add_separator(slide, 0.6, 1.05, 8.0)
    scenarios = [
        ("Scale = 0.5\u00d7",
         "Miniature world\nBaseline = 0.5 m\nBuilding height = 5 m",
         ACCENT3, 0.6),
        ("Scale = 1\u00d7 (true)",
         "Real world\nBaseline = 1.0 m\nBuilding height = 10 m",
         ACCENT, 4.85),
        ("Scale = 2\u00d7",
         "Giant world\nBaseline = 2.0 m\nBuilding height = 20 m",
         ACCENT2, 9.1),
    ]
    for title, desc, col, x in scenarios:
        add_rounded_rect(slide, x, 1.5, 3.8, 3.3, BG_MID)
        add_textbox(slide, x + 0.1, 1.6, 3.6, 0.5, title,
                    22, col, True, PP_ALIGN.CENTER)
        bld = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(x + 1.3), Inches(2.3),
            Inches(1.2), Inches(2.0))
        bld.fill.solid()
        bld.fill.fore_color.rgb = RGBColor(
            col[0] // 3, col[1] // 3, col[2] // 3)
        bld.line.fill.background()
        for cx_off in [0.3, 2.7]:
            cam = slide.shapes.add_shape(
                MSO_SHAPE.ISOSCELES_TRIANGLE,
                Inches(x + cx_off), Inches(3.2),
                Inches(0.4), Inches(0.3))
            cam.fill.solid(); cam.fill.fore_color.rgb = col
            cam.line.fill.background()
        add_textbox(slide, x + 0.1, 4.0, 3.6, 0.6, desc,
                    14, LIGHT_GRAY, False, PP_ALIGN.CENTER)
    add_textbox(slide, 0.6, 5.2, 12, 0.5,
                "All three reconstructions produce IDENTICAL 2D images!",
                24, YELLOW, True, PP_ALIGN.CENTER)
    add_textbox(slide, 0.6, 5.8, 12, 1.0,
                "The reprojection error is zero for all three. SfM has "
                "no way to distinguish between them.\nThe 3D shape is "
                "correct in all cases \u2014 only the absolute size differs.",
                18, LIGHT_GRAY, False, PP_ALIGN.CENTER)

    # ── SLIDE 9: Resolving Scale Ambiguity ──────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_textbox(slide, 0.6, 0.3, 12, 0.8,
                "How to Resolve Scale Ambiguity", 36, ACCENT, True)
    add_separator(slide, 0.6, 1.05, 5.5)
    methods = [
        ("Known Object Size",
         "Place an object of known dimensions in the scene.\n"
         "Measure it in the reconstruction \u2192 compute scale factor.",
         ACCENT3, 0.6, 1.4),
        ("Known Baseline",
         "If the distance between two camera positions is\n"
         "known (e.g., fixed stereo rig), scale is determined.",
         ACCENT3, 0.6, 2.8),
        ("IMU / Odometry",
         "Inertial Measurement Units provide metric acceleration.\n"
         "Visual-Inertial Odometry (VIO) fuses vision + IMU "
         "\u2192 metric scale.",
         ACCENT, 6.8, 1.4),
        ("Stereo Camera",
         "A calibrated stereo pair has a known baseline.\n"
         "Each stereo frame provides metric depth \u2192 absolute scale.",
         ACCENT, 6.8, 2.8),
        ("GPS / GNSS",
         "Geo-tagged images provide absolute positions.\n"
         "At least 2 positions fix the scale.",
         YELLOW, 0.6, 4.2),
        ("LiDAR Fusion",
         "LiDAR provides metric depth measurements.\n"
         "Fusing LiDAR with SfM resolves scale.",
         YELLOW, 6.8, 4.2),
    ]
    for title, desc, col, x, y in methods:
        add_rounded_rect(slide, x, y, 5.8, 1.2, BG_MID)
        add_textbox(slide, x + 0.2, y + 0.05, 5.4, 0.4,
                    title, 20, col, True)
        add_textbox(slide, x + 0.2, y + 0.5, 5.4, 0.7,
                    desc, 14, LIGHT_GRAY)
    add_textbox(slide, 0.6, 5.8, 12, 0.8,
                "Common thread:  every method introduces an external "
                "metric reference that pure images cannot provide.",
                18, YELLOW, True, PP_ALIGN.CENTER)

    # ── SLIDE 10: Summary ───────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_textbox(slide, 0.6, 0.3, 12, 0.8,
                "Summary: Scale Ambiguity in SfM", 36, ACCENT, True)
    add_separator(slide, 0.6, 1.05, 6.0)
    add_textbox(slide, 0.6, 1.3, 6, 0.5,
                "Key Takeaways", 24, ACCENT3, True)
    add_bullet_list(slide, 0.8, 1.9, 6, 4.5, [
        "SfM recovers 3D structure and camera motion from 2D "
        "images alone",
        "Perspective projection is invariant to global scaling "
        "of the scene",
        "Mathematically:  scaling X \u2192 \u03b1X  and  t \u2192 "
        "\u03b1t  produces identical images",
        "The essential matrix E = [t]\u2093R  is invariant to "
        "the magnitude of t",
        "No amount of monocular images can resolve this ambiguity",
        "Scale ambiguity is a fundamental property, not a "
        "limitation of algorithms",
        "External metric information (IMU, GPS, known distances) "
        "is required to recover absolute scale",
    ], 17, WHITE, Pt(10))
    add_rounded_rect(slide, 7.5, 1.5, 5.0, 2.2, BG_MID)
    add_textbox(slide, 7.7, 1.6, 4.6, 0.5,
                "The Fundamental Equation", 20, ACCENT, True,
                PP_ALIGN.CENTER)
    add_textbox(slide, 7.7, 2.2, 4.6, 1.0,
                "x  \u223c  K [R | t] X  \u223c  K [R | \u03b1t] "
                "(\u03b1X)", 22, YELLOW, True, PP_ALIGN.CENTER)
    add_textbox(slide, 7.7, 3.0, 4.6, 0.5,
                "\u2200 \u03b1 > 0:  same image!",
                20, ACCENT2, True, PP_ALIGN.CENTER)
    add_rounded_rect(slide, 7.5, 4.2, 5.0, 1.8, BG_MID)
    add_textbox(slide, 7.7, 4.3, 4.6, 0.4,
                "In Practice", 20, ACCENT, True, PP_ALIGN.CENTER)
    add_textbox(slide, 7.7, 4.8, 4.6, 1.1,
                "Set  ||t\u2081\u2082|| = 1  for initialization.\n"
                "The reconstruction is metrically\n"
                "consistent but not metrically correct\n"
                "until external scale is provided.",
                15, LIGHT_GRAY, False, PP_ALIGN.CENTER)
    add_textbox(slide, 0.6, 6.5, 12, 0.5,
                "\u201cStructure-from-Motion gives you the shape of "
                "the world, but not its size.\u201d",
                20, YELLOW, False, PP_ALIGN.CENTER)

    # Save
    out = os.path.join(os.path.dirname(__file__),
                       "Structure_from_Motion.pptx")
    prs.save(out)
    print(f"Saved: {out}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    build()
